"""Regression tests for document text extraction.

Guards the .pptx file-like bug: `_extract_pptx` must wrap raw bytes in
`io.BytesIO` because python-pptx's `Presentation()` needs a seekable file-like
object. Passing raw bytes raised "'bytes' object has no attribute 'seek'", which
caused every uploaded .pptx to be silently dropped (empty corpus -> HTTP 400 in
the AI course-generation flow).
"""
import io
from pptx import Presentation
from pptx.util import Inches
from services.document_service import DocumentService


def _pptx_bytes(text: str) -> bytes:
    prs = Presentation()
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1))
    tb.text_frame.text = text
    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()


def test_extract_text_from_pptx_returns_slide_text():
    data = _pptx_bytes("Track awareness is critical on rail sites")
    text = DocumentService.extract_text(data, "deck.pptx")
    assert "Track awareness is critical" in text


def test_extract_text_from_file_sync_pptx_by_content_type():
    data = _pptx_bytes("Emergency procedures and assembly points")
    ct = "application/vnd.openxmlformats-officedocument.presentationml.presentation"
    text = DocumentService.extract_text_from_file_sync(data, ct)
    assert "Emergency procedures" in text

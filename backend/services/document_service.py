"""
Document extraction service for various file formats.
Handles extraction of text from docx, pptx, and pdf files.
"""

import io
import logging
from typing import Optional
import fitz  # pymupdf — PDF extraction
from docx import Document as DocxDocument
from pptx import Presentation

logger = logging.getLogger(__name__)


class DocumentService:
    """Service for extracting text from documents."""

    @staticmethod
    def extract_text_from_file_sync(file_bytes: bytes, content_type: str) -> str:
        """
        Extract text from document file bytes, using content-type to determine format.

        Args:
            file_bytes: Raw file bytes
            content_type: MIME type (e.g. 'application/pdf', 'application/vnd.openxmlformats...')

        Returns:
            Extracted text content
        """
        ct = content_type.lower()
        if "pdf" in ct:
            return DocumentService._extract_pdf(file_bytes)
        elif "wordprocessingml" in ct or "docx" in ct or "msword" in ct:
            return DocumentService._extract_docx(file_bytes)
        elif "presentationml" in ct or "pptx" in ct or "powerpoint" in ct:
            return DocumentService._extract_pptx(file_bytes)
        else:
            # Fallback: try PDF, then DOCX
            try:
                return DocumentService._extract_pdf(file_bytes)
            except Exception:
                return DocumentService._extract_docx(file_bytes)

    @staticmethod
    def extract_text(file_bytes: bytes, filename: str) -> str:
        """
        Extract text from document file bytes.

        Args:
            file_bytes: Raw file bytes
            filename: Original filename (used to determine file type)

        Returns:
            Extracted text content

        Raises:
            ValueError: If file type is not supported
        """
        filename_lower = filename.lower()

        if filename_lower.endswith(".docx"):
            return DocumentService._extract_docx(file_bytes)
        elif filename_lower.endswith(".pptx"):
            return DocumentService._extract_pptx(file_bytes)
        elif filename_lower.endswith(".pdf"):
            return DocumentService._extract_pdf(file_bytes)
        else:
            raise ValueError(
                f"Unsupported file type: {filename}. "
                "Supported types: .docx, .pptx, .pdf"
            )

    @staticmethod
    def build_corpus(texts: list[str], max_chars: int = 60000) -> str:
        """Merge per-file extracted texts into a single capped corpus.

        Each file gets an equal share of the budget so one large file can't
        crowd out the others; the joined result is hard-capped at max_chars.
        """
        non_empty = [t for t in texts if t and t.strip()]
        if not non_empty:
            return ""
        per_file = max(1, max_chars // len(non_empty))
        parts = [t[:per_file] for t in non_empty]
        corpus = "\n\n---\n\n".join(parts)
        return corpus[:max_chars]

    @staticmethod
    def _extract_docx(file_bytes: bytes) -> str:
        """
        Extract text from DOCX file.

        Args:
            file_bytes: Raw file bytes

        Returns:
            Extracted text with structure preserved
        """
        try:
            doc = DocxDocument(io.BytesIO(file_bytes))  # BytesIO wraps raw bytes
            content_parts = []

            for para in doc.paragraphs:
                if para.text.strip():
                    content_parts.append(para.text)

            # Extract text from tables if present
            for table in doc.tables:
                table_text = []
                for row in table.rows:
                    row_text = [cell.text for cell in row.cells]
                    table_text.append(" | ".join(row_text))
                if table_text:
                    content_parts.append("\n".join(table_text))

            extracted_text = "\n".join(content_parts)
            logger.info(f"Extracted text from DOCX: {len(extracted_text)} characters")
            return extracted_text

        except Exception as e:
            logger.error(f"Error extracting DOCX: {str(e)}")
            raise ValueError(f"Failed to extract text from DOCX file: {str(e)}")

    @staticmethod
    def _extract_pptx(file_bytes: bytes) -> str:
        """
        Extract text from PPTX file (slide by slide, including notes).

        Args:
            file_bytes: Raw file bytes

        Returns:
            Extracted text organized by slide
        """
        try:
            presentation = Presentation(io.BytesIO(file_bytes))  # BytesIO: python-pptx needs a seekable file-like object
            content_parts = []

            for slide_idx, slide in enumerate(presentation.slides, 1):
                content_parts.append(f"--- Slide {slide_idx} ---")

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        content_parts.append(shape.text)

                # Extract speaker notes if available
                if slide.has_notes_slide:
                    notes_text = slide.notes_slide.notes_text_frame.text.strip()
                    if notes_text:
                        content_parts.append(f"Notes: {notes_text}")

            extracted_text = "\n".join(content_parts)
            logger.info(f"Extracted text from PPTX: {len(extracted_text)} characters")
            return extracted_text

        except Exception as e:
            logger.error(f"Error extracting PPTX: {str(e)}")
            raise ValueError(f"Failed to extract text from PPTX file: {str(e)}")

    @staticmethod
    def _extract_pdf(file_bytes: bytes) -> str:
        """
        Extract text from PDF file using PyMuPDF (fitz).

        Args:
            file_bytes: Raw file bytes

        Returns:
            Extracted text (best effort)

        Note:
            Uses fitz.open(stream=...) for reliable text extraction from text-based PDFs.
            Image-only / scanned PDFs will produce minimal text; a warning is logged.
        """
        try:
            doc = fitz.open(stream=file_bytes, filetype="pdf")
            content_parts = []
            for page in doc:
                text = page.get_text()
                if text.strip():
                    content_parts.append(text)
            doc.close()
            extracted_text = "\n".join(content_parts)
            if len(extracted_text.strip()) < 100:
                logger.warning("PDF extraction produced minimal text — may be scanned/image-only")
            logger.info(f"Extracted {len(extracted_text)} characters from PDF via PyMuPDF")
            return extracted_text

        except Exception as e:
            logger.error(f"Error extracting PDF: {str(e)}")
            raise ValueError(f"Failed to extract text from PDF: {str(e)}")

"""
Slide generation service.
Converts generated course content into a PowerPoint presentation.
No Claude call needed — uses the existing course structure directly.
"""

import io
import logging
import textwrap
from typing import Any, Dict, List

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

logger = logging.getLogger(__name__)

# Brand colours
DARK_BG = RGBColor(0x1F, 0x29, 0x37)       # #1F2937 — sidebar navy
ACCENT = RGBColor(0x63, 0x66, 0xF1)         # #6366F1 — indigo
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF3, 0xF4, 0xF6)
DARK_TEXT = RGBColor(0x1F, 0x29, 0x37)
MID_TEXT = RGBColor(0x6B, 0x72, 0x80)


def _set_bg(slide, color: RGBColor):
    """Set solid background colour on a slide."""
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, text, font_size, bold=False,
                 color=WHITE, align=PP_ALIGN.LEFT, wrap=True):
    """Add a textbox to a slide and return the text frame."""
    txBox = slide.shapes.add_textbox(
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color
    return tf


def _key_points(lesson: Dict[str, Any], max_points: int = 4) -> List[str]:
    """Extract key points from a lesson — outcomes first, then content sentences."""
    outcomes = lesson.get("learning_outcomes", [])
    if outcomes:
        return [str(o) for o in outcomes[:max_points]]

    content = str(lesson.get("content", ""))
    sentences = [s.strip() for s in content.replace("\n", " ").split(".") if s.strip()]
    return sentences[:max_points]


class SlideService:
    """Generates PPTX presentations from course content."""

    def generate_slides(self, course_content: Dict[str, Any]) -> bytes:
        """
        Build a PPTX presentation from course content.

        Args:
            course_content: The course content dict (stored in Course.content)

        Returns:
            PPTX file as bytes
        """
        prs = Presentation()
        prs.slide_width = Inches(13.33)
        prs.slide_height = Inches(7.5)

        title = course_content.get("title", "Untitled Course")
        description = course_content.get("description", "")
        modules = course_content.get("modules", [])
        objectives = course_content.get("objectives", [])

        logger.info(f"Building slides for: {title}")

        # Slide 1 — Title
        self._add_title_slide(prs, title, description)

        # Slide 2 — Objectives (if present)
        if objectives:
            self._add_objectives_slide(prs, objectives)

        # Module + lesson slides
        for module in modules:
            self._add_module_divider(prs, module)
            for lesson in module.get("lessons", []):
                self._add_lesson_slide(prs, module, lesson)

        # Final slide
        self._add_closing_slide(prs, title)

        buffer = io.BytesIO()
        prs.save(buffer)
        buffer.seek(0)
        logger.info(f"Slides built for: {title} ({len(prs.slides)} slides)")
        return buffer.read()

    def _add_title_slide(self, prs: Presentation, title: str, description: str):
        slide_layout = prs.slide_layouts[6]  # blank
        slide = prs.slides.add_slide(slide_layout)
        _set_bg(slide, DARK_BG)

        # Accent bar
        bar = slide.shapes.add_shape(1, Inches(0), Inches(3.4), Inches(13.33), Inches(0.06))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        _add_textbox(slide, 0.8, 1.5, 11.5, 1.6, title,
                     font_size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        if description:
            desc = description[:200] + ("..." if len(description) > 200 else "")
            _add_textbox(slide, 1.5, 3.8, 10, 1.2, desc,
                         font_size=18, color=RGBColor(0xD1, 0xD5, 0xDB), align=PP_ALIGN.CENTER)

        _add_textbox(slide, 0.8, 6.4, 11.5, 0.6, "Learning & Development",
                     font_size=12, color=ACCENT, align=PP_ALIGN.CENTER)

    def _add_objectives_slide(self, prs: Presentation, objectives: List[str]):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        _set_bg(slide, LIGHT_GRAY)

        _add_textbox(slide, 0.6, 0.3, 12, 0.7, "Learning Objectives",
                     font_size=28, bold=True, color=DARK_BG)

        # Accent underline
        bar = slide.shapes.add_shape(1, Inches(0.6), Inches(1.1), Inches(2.5), Inches(0.05))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        txBox = slide.shapes.add_textbox(Inches(0.6), Inches(1.4), Inches(12), Inches(5.5))
        tf = txBox.text_frame
        tf.word_wrap = True

        for i, obj in enumerate(objectives[:6]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.space_before = Pt(6)
            run = p.add_run()
            run.text = f"✓  {obj}"
            run.font.size = Pt(18)
            run.font.color.rgb = DARK_TEXT

    def _add_module_divider(self, prs: Presentation, module: Dict[str, Any]):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        _set_bg(slide, ACCENT)

        mod_num = f"Module {module.get('id', '')}"
        _add_textbox(slide, 0.8, 1.8, 11.5, 0.6, mod_num,
                     font_size=18, color=WHITE, align=PP_ALIGN.CENTER)

        _add_textbox(slide, 0.8, 2.5, 11.5, 1.6, module.get("title", ""),
                     font_size=40, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        if module.get("description"):
            _add_textbox(slide, 1.5, 4.3, 10, 1.2, module["description"][:180],
                         font_size=16, color=RGBColor(0xE0, 0xE7, 0xFF), align=PP_ALIGN.CENTER)

    def _add_lesson_slide(self, prs: Presentation, module: Dict[str, Any], lesson: Dict[str, Any]):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        _set_bg(slide, WHITE)

        # Top accent strip
        bar = slide.shapes.add_shape(1, Inches(0), Inches(0), Inches(13.33), Inches(0.08))
        bar.fill.solid()
        bar.fill.fore_color.rgb = ACCENT
        bar.line.fill.background()

        # Module label
        _add_textbox(slide, 0.5, 0.2, 12, 0.4,
                     f"Module {module.get('id', '')} — {module.get('title', '')}",
                     font_size=11, color=MID_TEXT)

        # Lesson title
        _add_textbox(slide, 0.5, 0.7, 12, 0.9, lesson.get("title", ""),
                     font_size=30, bold=True, color=DARK_BG)

        # Accent underline
        bar2 = slide.shapes.add_shape(1, Inches(0.5), Inches(1.7), Inches(1.5), Inches(0.05))
        bar2.fill.solid()
        bar2.fill.fore_color.rgb = ACCENT
        bar2.line.fill.background()

        # Key points
        points = _key_points(lesson)
        if points:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(2.0), Inches(12.3), Inches(5.0))
            tf = txBox.text_frame
            tf.word_wrap = True

            for i, point in enumerate(points):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.space_before = Pt(10)
                run = p.add_run()
                # Truncate long points
                text = point if len(point) < 160 else point[:157] + "..."
                run.text = f"▸  {text}"
                run.font.size = Pt(20)
                run.font.color.rgb = DARK_TEXT

        # Duration badge
        duration = lesson.get("estimated_duration_minutes")
        if duration:
            _add_textbox(slide, 10.8, 6.9, 2.3, 0.4, f"⏱ {duration} min",
                         font_size=11, color=MID_TEXT, align=PP_ALIGN.RIGHT)

    def _add_closing_slide(self, prs: Presentation, course_title: str):
        slide_layout = prs.slide_layouts[6]
        slide = prs.slides.add_slide(slide_layout)
        _set_bg(slide, DARK_BG)

        _add_textbox(slide, 0.8, 2.0, 11.5, 1.2, "Thank You",
                     font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

        _add_textbox(slide, 0.8, 3.4, 11.5, 0.8, "Questions & Discussion",
                     font_size=24, color=ACCENT, align=PP_ALIGN.CENTER)

        _add_textbox(slide, 0.8, 5.6, 11.5, 0.5, course_title,
                     font_size=14, color=RGBColor(0x9C, 0xA3, 0xAF), align=PP_ALIGN.CENTER)
